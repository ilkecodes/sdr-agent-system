"""Canonical converter: heterogeneous docs -> loss-aware GFM Markdown + RAG-optimized chunks

Usage (high-level):
  from app.convert import convert_file
  convert_file(source_uri='/abs/path/to/file.pdf', file_bytes=open(...,'rb').read(), out_dir='out')

Produces two files (by default in out_dir):
  {doc_stem}.md
  {doc_stem}.chunks.jsonl

This implementation follows the project's ruleset delivered by the user; it is a pragmatic, dependency-friendly prototype
that favors deterministic behavior and clear metadata. It covers txt, pdf, docx, pptx, xlsx/csv, json and html where
libraries are available. It uses tiktoken for token estimation when available, otherwise falls back to a word-based
heuristic.
"""

from __future__ import annotations

import os
import re
import io
import json
import math
import hashlib
import datetime
from typing import List, Dict, Any, Tuple, Optional

try:
    import tiktoken
    def count_tokens(text: str, enc_name: str = 'cl100k_base') -> int:
        try:
            enc = tiktoken.get_encoding(enc_name)
            return len(enc.encode(text))
        except Exception:
            return len(text.split())
except Exception:
    def count_tokens(text: str, enc_name: str = 'cl100k_base') -> int:
        # fallback approximate: words ~ tokens
        return max(1, len(text.split()))

try:
    # prefer pypdf (successor to PyPDF2)
    from pypdf import PdfReader
except Exception:
    try:
        from PyPDF2 import PdfReader
    except Exception:
        PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pandas as pd
except Exception:
    pd = None

from html import escape as html_escape

MAX_TABLE_WIDTH = 6  # columns for GFM before falling back to HTML
TARGET_TOKENS = 1200
OVERLAP_TOKENS = 200


def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def deterministic_chunk_id(source_uri: str, heading_path: str, range_id: str) -> str:
    s = f"{source_uri}|{heading_path}|{range_id}"
    return hashlib.sha256(s.encode('utf-8')).hexdigest()[:16]


def safe_text(b: bytes) -> str:
    for enc in ('utf-8', 'latin-1', 'utf-16'):
        try:
            return b.decode(enc)
        except Exception:
            continue
    return b.decode('latin-1', errors='ignore')


def simple_keywords(text: str, max_k: int = 8) -> List[str]:
    # naive keyword extraction: freq of words excluding stopwords and short tokens
    stop = set(["the","and","of","in","to","a","is","for","on","with","that","by","as","an","are","be","this","it"])
    words = re.findall(r"[A-Za-z0-9%$\-]+", text.lower())
    freq: Dict[str, int] = {}
    for w in words:
        if w in stop or len(w) < 3:
            continue
        freq[w] = freq.get(w, 0) + 1
    items = sorted(freq.items(), key=lambda x: (-x[1], x[0]))
    return [k for k, _ in items[:max_k]]


def md_table_from_dataframe(df) -> Tuple[str, bool]:
    """Return (markdown_table, used_html_fallback)
    GFM table if columns <= MAX_TABLE_WIDTH, otherwise HTML fallback with colspan/rowspan not handled here (best-effort).
    """
    cols = list(df.columns)
    if len(cols) <= MAX_TABLE_WIDTH:
        # build GFM with aligned pipes
        headers = "| " + " | ".join(str(c) for c in cols) + " |\n"
        sep = "| " + " | ".join(['---'] * len(cols)) + " |\n"
        rows = []
        for _, r in df.iterrows():
            cells = [escape_pipe(str(r[c])) for c in cols]
            rows.append("| " + " | ".join(cells) + " |\n")
        return headers + sep + ''.join(rows), False
    else:
        # fallback to HTML table
        html = ["<table>\n<thead>\n<tr>"]
        for c in cols:
            html.append(f"<th>{html_escape(str(c))}</th>")
        html.append("</tr>\n</thead>\n<tbody>\n")
        for _, r in df.iterrows():
            html.append("<tr>")
            for c in cols:
                html.append(f"<td>{html_escape(str(r[c]))}</td>")
            html.append("</tr>\n")
        html.append("</tbody>\n</table>\n")
        return ''.join(html), True


def escape_pipe(text: str) -> str:
    return text.replace('|', '\\|')


def front_matter(metadata: Dict[str, Any]) -> str:
    # emit a compact YAML-like front matter; avoid adding a pyyaml dependency
    lines = ['---']
    for k, v in metadata.items():
        if isinstance(v, (dict, list)):
            j = json.dumps(v, ensure_ascii=False)
            lines.append(f"{k}: '{j}'")
        else:
            lines.append(f"{k}: '{str(v)}'")
    lines.append('---\n')
    return "\n".join(lines)


def chunk_markdown(blocks: List[str], source_uri: str, heading_paths: List[str], target_tokens: int = TARGET_TOKENS, overlap_tokens: int = OVERLAP_TOKENS) -> List[Dict[str, Any]]:
    """Take markdown blocks (atomic units) and produce list of chunk dicts with deterministic ids and metadata.
    heading_paths is parallel to blocks and gives the heading_path for each block.
    """
    chunks: List[Dict[str, Any]] = []
    cur: List[str] = []
    cur_tokens = 0
    cur_heading = ''
    byte_cursor = 0

    def flush(reason: str):
        nonlocal cur, cur_tokens, cur_heading, byte_cursor
        if not cur:
            return
        text = '\n\n'.join(cur)
        tok = count_tokens(text)
        # deterministically assign range id using byte_cursor and tok
        range_id = f"{byte_cursor}-{byte_cursor + tok}"
        chunk_id = deterministic_chunk_id(source_uri, cur_heading, range_id)
        summary = first_sentence(text)
        keywords = simple_keywords(text)
        meta = {
            'heading_path': cur_heading,
            'byte_range': [byte_cursor, byte_cursor + tok],
            'tokens': tok,
        }
        chunks.append({
            'chunk_id': chunk_id,
            'text': text,
            'summary': summary,
            'keywords': keywords,
            'metadata': meta,
        })
        byte_cursor += tok
        cur = []
        cur_tokens = 0

    for b, hp in zip(blocks, heading_paths):
        toks = count_tokens(b)
        # reset current heading when the heading path changes
        if hp != cur_heading:
            # flush current chunk at section break
            flush('section_break')
            cur_heading = hp

        # atomic blocks should start their own chunk if they'd overflow
        is_atomic = b.strip().startswith('|') or b.strip().startswith('```') or b.strip().startswith('<table')

        if is_atomic and cur and (cur_tokens + toks > target_tokens):
            flush('atomic_boundary')

        cur.append(b)
        cur_tokens += toks

        if cur_tokens >= target_tokens:
            flush('size_limit')
            # overlap handling: copy tail tokens from last chunk
            if overlap_tokens > 0 and len(chunks) > 0:
                tail = take_tail_tokens(chunks[-1]['text'], overlap_tokens)
                cur = [tail]
                cur_tokens = count_tokens(tail)

    if cur:
        flush('final')

    return chunks


def take_tail_tokens(text: str, n_tokens: int) -> str:
    toks = text.split()
    if len(toks) <= n_tokens:
        return text
    return ' '.join(toks[-n_tokens:])


def first_sentence(text: str) -> str:
    m = re.search(r"(.+?[\.\!\?])(\s|$)", text.strip())
    if m:
        s = m.group(1).strip()
        return s if len(s) <= 300 else s[:297] + '...'
    return text.strip()[:300]


def convert_file(source_uri: str, file_bytes: Optional[bytes] = None, mime: Optional[str] = None, doc_lang_hint: str = 'auto', out_dir: Optional[str] = None) -> Dict[str, Any]:
    """Main entry point. Returns dict with paths written.
    If file_bytes is None, source_uri is treated as a local path and file is read.
    """
    if file_bytes is None:
        with open(source_uri, 'rb') as f:
            file_bytes = f.read()

    doc_stem = os.path.splitext(os.path.basename(source_uri))[0]
    out_dir = out_dir or os.getcwd()
    os.makedirs(out_dir, exist_ok=True)

    checksum = sha256_hex(file_bytes)
    extracted_blocks: List[str] = []
    heading_paths: List[str] = []
    metadata_top: Dict[str, Any] = {
        'source_uri': os.path.abspath(source_uri),
        'checksum_sha256': checksum,
        'content_type': None,
        'mime': mime or 'application/octet-stream',
        'doc_language': doc_lang_hint,
        'extraction_engine': 'convert.py:prototype',
        'extraction_version': '0.1',
        # timezone-aware UTC timestamp
        'created_at': datetime.datetime.now(datetime.timezone.utc).isoformat(),
    }

    lower = source_uri.lower()
    try:
        if lower.endswith('.txt'):
            text = safe_text(file_bytes)
            for i, para in enumerate([p.strip() for p in text.split('\n\n') if p.strip()]):
                heading = 'Root' if i == 0 else 'Root'
                extracted_blocks.append(para)
                heading_paths.append(heading)
            metadata_top['content_type'] = 'txt'

        elif lower.endswith('.pdf') and PdfReader is not None:
            metadata_top['content_type'] = 'pdf'
            reader = PdfReader(io.BytesIO(file_bytes))
            for pnum, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ''
                # naive paragraph split
                paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
                heading = f"Page {pnum}"
                for para in paras:
                    extracted_blocks.append(para)
                    heading_paths.append(heading)

        elif lower.endswith('.docx') and docx is not None:
            metadata_top['content_type'] = 'docx'
            doc = docx.Document(io.BytesIO(file_bytes))
            cur_heading = []
            for para in doc.paragraphs:
                style = para.style.name if para.style is not None else ''
                text = para.text.strip()
                if not text:
                    continue
                if style.lower().startswith('heading'):
                    # style names are usually 'Heading 1'..'Heading 6'
                    level = re.search(r"(\d+)", style)
                    level = int(level.group(1)) if level else 2
                    cur_heading = cur_heading[:level-1] + [text]
                    extracted_blocks.append('#' * level + ' ' + text)
                    heading_paths.append(' > '.join(cur_heading))
                else:
                    extracted_blocks.append(text)
                    heading_paths.append(' > '.join(cur_heading) or 'Root')

        elif (lower.endswith('.pptx') or lower.endswith('.ppt')) and Presentation is not None:
            metadata_top['content_type'] = 'pptx'
            pres = Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(pres.slides, start=1):
                title = None
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        txt = shape.text.strip()
                        if not title:
                            title = txt
                            break
                title = title or f"Slide {i}"
                heading = f"Slide: {title}"
                extracted_blocks.append(f"# Slide: {title}")
                heading_paths.append(heading)
                # bullets -> nested lists
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for p in shape.text_frame.paragraphs:
                        level = p.level if hasattr(p, 'level') else 0
                        prefix = ('  ' * level) + ('- ' if p.level is not None and p.level >= 0 else '- ')
                        extracted_blocks.append(prefix + p.text.strip())
                        heading_paths.append(heading)
                # speaker notes if available -> details block
                notes = ''
                if slide.has_notes_slide:
                    try:
                        notes = '\n'.join([p.text for p in slide.notes_slide.notes_text_frame.paragraphs])
                    except Exception:
                        notes = ''
                if notes:
                    extracted_blocks.append(f"<details><summary>Speaker notes</summary>\n\n{notes}\n\n</details>")
                    heading_paths.append(heading)

        elif (lower.endswith('.xls') or lower.endswith('.xlsx') or lower.endswith('.csv')) and pd is not None:
            metadata_top['content_type'] = 'xlsx' if 'xls' in lower else 'csv'
            if lower.endswith('.csv'):
                df = pd.read_csv(io.BytesIO(file_bytes))
                md, html_fallback = md_table_from_dataframe(df)
                caption = f"*Table: {doc_stem} (csv)*\n"
                if html_fallback:
                    extracted_blocks.append(caption + md)
                else:
                    extracted_blocks.append(caption + md)
                heading_paths.append('Sheet: default')
            else:
                xls = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
                for sheetname, df in xls.items():
                    caption = f"*Table: {sheetname}*\n"
                    md, html_fallback = md_table_from_dataframe(df)
                    if html_fallback:
                        extracted_blocks.append(caption + md)
                    else:
                        extracted_blocks.append(caption + md)
                    heading_paths.append(f"Sheet: {sheetname}")

        elif lower.endswith('.json'):
            metadata_top['content_type'] = 'json'
            text = safe_text(file_bytes)
            try:
                obj = json.loads(text)
                pretty = json.dumps(obj, ensure_ascii=False, indent=2)
                # If large array, summarize and sample
                if isinstance(obj, (list, dict)) and len(pretty) > 2000:
                    summary = f"JSON object with top-level type {type(obj).__name__} and {len(obj) if hasattr(obj,'__len__') else 'unknown'} items"
                    extracted_blocks.append(summary)
                    heading_paths.append('Root')
                    sample = json.dumps(obj[:5] if isinstance(obj, list) else list(obj.items())[:5], ensure_ascii=False, indent=2)
                    extracted_blocks.append('```json\n' + sample + '\n```')
                    heading_paths.append('Root')
                else:
                    extracted_blocks.append('```json\n' + pretty + '\n```')
                    heading_paths.append('Root')
            except Exception:
                extracted_blocks.append('```json\n' + safe_text(file_bytes)[:2000] + '\n```')
                heading_paths.append('Root')

        elif lower.endswith('.html') or lower.endswith('.htm'):
            metadata_top['content_type'] = 'html'
            text = safe_text(file_bytes)
            # naive: extract H1-H3 via regex
            for m in re.finditer(r"<h([1-6])[^>]*>(.*?)</h\1>", text, flags=re.I|re.S):
                level = int(m.group(1))
                content = re.sub(r"<[^>]+>", '', m.group(2)).strip()
                extracted_blocks.append('#' * level + ' ' + content)
                heading_paths.append(content)
            # fallback: extract main paragraphs
            paras = re.split(r"</p>\s*", text, flags=re.I)
            for p in paras[:20]:
                s = re.sub(r"<[^>]+>", '', p).strip()
                if s:
                    extracted_blocks.append(s)
                    heading_paths.append('Root')

        else:
            # unknown type: try text
            metadata_top['content_type'] = 'unknown'
            extracted_blocks.append(safe_text(file_bytes)[:2000])
            heading_paths.append('Root')

    except Exception as e:
        # On extractor failure, emit partial md and marker
        extracted_blocks.append(f"<!-- extraction error: {str(e)} -->")
        heading_paths.append('Root')

    # Build canonical Markdown: front matter + content
    md_meta = {
        **metadata_top,
        'title': doc_stem,
    }
    md_lines = [front_matter(md_meta)]
    for b in extracted_blocks:
        md_lines.append(b)
        md_lines.append('\n')

    md_text = '\n'.join(md_lines)

    # Chunking
    chunks = chunk_markdown(extracted_blocks, metadata_top['source_uri'], heading_paths)

    # Attach richer metadata to each chunk and normalize
    for c in chunks:
        c['metadata'].update({
            'source_uri': metadata_top['source_uri'],
            'checksum_sha256': metadata_top['checksum_sha256'],
            'content_type': metadata_top.get('content_type'),
            'extraction_engine': metadata_top['extraction_engine'],
            'extraction_version': metadata_top['extraction_version'],
            'doc_language': metadata_top['doc_language'],
        })

    md_path = os.path.join(out_dir, f"{doc_stem}.md")
    chunks_path = os.path.join(out_dir, f"{doc_stem}.chunks.jsonl")

    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(md_text)

    with open(chunks_path, 'w', encoding='utf-8') as f:
        for c in chunks:
            out = {
                'chunk_id': c['chunk_id'],
                'text': c['text'],
                'summary': c['summary'],
                'keywords': c['keywords'],
                'metadata': c['metadata'],
            }
            f.write(json.dumps(out, ensure_ascii=False) + '\n')

    return {'md_path': md_path, 'chunks_path': chunks_path, 'n_chunks': len(chunks)}


if __name__ == '__main__':
    import argparse
    p = argparse.ArgumentParser(description='Convert file -> markdown + RAG chunks')
    p.add_argument('source', help='path to file')
    p.add_argument('--out', help='output directory', default='out')
    args = p.parse_args()
    res = convert_file(args.source, out_dir=args.out)
    print('Wrote:', res)
